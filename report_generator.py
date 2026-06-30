import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

from branding import COLORS, get_rgb
from data_parser import STATS, ATTACK_STATS, DEFENSE_STATS, H2H_STATS
from insights import generate_insights

try:
    from pptx.enum.shapes import MSO_SHAPE
    RECT = MSO_SHAPE.RECTANGLE
    RRECT = MSO_SHAPE.ROUNDED_RECTANGLE
except Exception:
    from pptx.util import Emu
    RECT = 1   # fallback integer id
    RRECT = 5


def safe_float(val):
    try:
        return float(val)
    except (ValueError, TypeError):
        return 0.0


# ── PRESENTATION SETUP ──────────────────────────────────────────────────────
def create_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


# ── SHARED HELPERS ───────────────────────────────────────────────────────────
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_shape(slide, left, top, width, height, fill_hex, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = get_rgb(COLORS[fill_hex])
    if line:
        shape.line.color.rgb = get_rgb(COLORS[fill_hex])
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              size=12, bold=False, color="BLACK",
              align=PP_ALIGN.LEFT, italic=False, font="Arial"):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = get_rgb(COLORS[color])
    p.font.name = font
    p.alignment = align
    return box


def _add_slide_header(slide, title_text, subtitle=None):
    """Standard black banner header used on every content slide."""
    # Banner
    _add_shape(slide, 0, 0, 13.33, 0.82, "BLACK")
    # Gold stripe
    _add_shape(slide, 0, 0.82, 13.33, 0.05, "GOLD")
    # Title
    _add_text(slide, title_text,
              0.45, 0.12, 12.5, 0.62,
              size=22, bold=True, color="GOLD", font="Arial Black")
    if subtitle:
        _add_text(slide, subtitle,
                  0.45, 0.58, 12.5, 0.3,
                  size=10, bold=False, color="SILVER")


def _add_slide_footer(slide, data):
    """Consistent footer on every content slide."""
    footer_text = (
        f"Brooklyn FC  |  {data['opponent_name'].upper()}  |  "
        f"{data['match_date']}  |  {data['competition']}"
    )
    _add_text(slide, footer_text,
              0, 7.18, 13.33, 0.28,
              size=7.5, italic=True, color="SILVER",
              align=PP_ALIGN.CENTER)


def _set_cell(cell, text, size=11, bold=False,
              color="BLACK", align=PP_ALIGN.LEFT,
              bg_hex=None):
    if bg_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_rgb(COLORS[bg_hex])
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = get_rgb(COLORS[color])
    p.alignment = align


# ── CHART ENGINE ─────────────────────────────────────────────────────────────
def _bar_chart(values, labels, colors, title, width=5.2, height=3.8, dpi=180):
    """
    Render a clean bar chart and return an in-memory PNG BytesIO stream.
    colors — list of hex strings (without #).
    """
    fig, ax = plt.subplots(figsize=(width, height), dpi=dpi)
    fig.patch.set_facecolor("#" + COLORS["WHITE"])
    ax.set_facecolor("#" + COLORS["WHITE"])

    x = np.arange(len(labels))
    bar_colors = ["#" + c for c in colors]
    bars = ax.bar(x, values, color=bar_colors, width=0.45, edgecolor="none", zorder=3)

    ax.set_title(title.upper(), fontsize=9, fontweight="bold",
                 pad=10, color="#" + COLORS["BLACK"])
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=8.5, color="#" + COLORS["BLACK"])
    ax.tick_params(axis="y", colors="#" + COLORS["SILVER"], labelsize=7.5)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.spines["bottom"].set_color("#" + COLORS["SILVER"])
    ax.yaxis.grid(True, linestyle="--", alpha=0.45, color="#" + COLORS["GRID"], zorder=0)
    ax.set_axisbelow(True)

    for bar in bars:
        h = bar.get_height()
        label_val = f"{h:.2g}" if h != int(h) else f"{int(h)}"
        ax.annotate(label_val,
                    xy=(bar.get_x() + bar.get_width() / 2, h),
                    xytext=(0, 4), textcoords="offset points",
                    ha="center", va="bottom",
                    fontsize=9.5, fontweight="bold",
                    color="#" + COLORS["BLACK"])


    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor(), edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


def _h_bar_chart(bkfc_vals, opp_vals, labels, title, width=10.5, height=4.8, dpi=160):
    """
    Horizontal grouped bar chart for head-to-head comparison.
    """
    n = len(labels)
    y  = np.arange(n)
    h  = 0.35

    fig, ax = plt.subplots(figsize=(width, height), dpi=dpi)
    fig.patch.set_facecolor("#" + COLORS["WHITE"])
    ax.set_facecolor("#" + COLORS["WHITE"])

    b1 = ax.barh(y + h/2, bkfc_vals, h, color="#" + COLORS["GOLD"],
                 label="BKFC", edgecolor="none")
    b2 = ax.barh(y - h/2, opp_vals, h, color="#" + COLORS["DARK_GRAY"],
                 label="Opponent", edgecolor="none")

    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=8.5, color="#" + COLORS["BLACK"])
    ax.tick_params(axis="x", colors="#" + COLORS["SILVER"], labelsize=7.5)
    ax.set_title(title.upper(), fontsize=10, fontweight="bold",
                 pad=10, color="#" + COLORS["BLACK"])

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_color("#" + COLORS["SILVER"])
    ax.spines["left"].set_visible(False)
    ax.xaxis.grid(True, linestyle="--", alpha=0.4, color="#" + COLORS["GRID"], zorder=0)
    ax.set_axisbelow(True)

    for bar in b1:
        w = bar.get_width()
        ax.text(w + max(bkfc_vals + opp_vals) * 0.01, bar.get_y() + bar.get_height()/2,
                f"{w:.2g}", va="center", ha="left", fontsize=7.5, fontweight="bold",
                color="#" + COLORS["BLACK"])
    for bar in b2:
        w = bar.get_width()
        ax.text(w + max(bkfc_vals + opp_vals) * 0.01, bar.get_y() + bar.get_height()/2,
                f"{w:.2g}", va="center", ha="left", fontsize=7.5,
                color="#" + COLORS["DARK_GRAY"])

    ax.legend(fontsize=8, frameon=False, loc="lower right")
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor(), edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── SLIDE 1: TITLE ───────────────────────────────────────────────────────────
def add_title_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])

    # Gold top stripe
    _add_shape(slide, 0, 0, 13.33, 0.1, "GOLD")
    # Gold bottom stripe
    _add_shape(slide, 0, 7.4, 13.33, 0.1, "GOLD")

    # Club name
    _add_text(slide, "BROOKLYN FC",
              1, 1.5, 11.33, 1.4,
              size=62, bold=True, color="GOLD",
              align=PP_ALIGN.CENTER, font="Arial Black")

    # Sub heading
    _add_text(slide, "MATCH ANALYSIS REPORT",
              1, 2.85, 11.33, 0.55,
              size=16, color="WHITE",
              align=PP_ALIGN.CENTER)

    # Score badge background
    score_bg = slide.shapes.add_shape(
        1, Inches(5.16), Inches(3.75), Inches(3.0), Inches(0.75)
    )
    score_bg.fill.solid()
    score_bg.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    score_bg.line.fill.background()

    # Score text (on top of badge)
    _add_text(slide, data["score"],
              5.16, 3.75, 3.0, 0.75,
              size=34, bold=True, color="BLACK",
              align=PP_ALIGN.CENTER, font="Arial Black")

    # Vs line
    _add_text(slide, f"BKFC  vs  {data['opponent_name'].upper()}",
              1, 4.75, 11.33, 0.65,
              size=22, bold=True, color="WHITE",
              align=PP_ALIGN.CENTER)

    # Metadata
    _add_text(slide,
              f"{data['match_date']}   |   {data['competition']}",
              1, 5.45, 11.33, 0.4,
              size=12, color="SILVER",
              align=PP_ALIGN.CENTER)


# ── SLIDE 2: MATCH PERFORMANCE OVERVIEW ─────────────────────────────────────
def add_summary_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, "MATCH PERFORMANCE OVERVIEW")
    _add_slide_footer(slide, data)

    opp = data["opponent_name"].upper()
    summary_metrics = [
        ("Goals",           6),
        ("xG (Exp. Goals)", 7),
        ("Shots",           8),
        ("Shots on Target", 9),
        ("Possession %",    14),
        ("Pass Accuracy %", 13),
        ("PPDA",            108),
    ]

    tbl = slide.shapes.add_table(
        rows=len(summary_metrics) + 1, cols=3,
        left=Inches(1.5), top=Inches(1.1),
        width=Inches(10.33), height=Inches(5.7)
    ).table

    # Header row
    for ci, txt in enumerate(["METRIC", "BROOKLYN FC", opp]):
        _set_cell(tbl.cell(0, ci), txt,
                  size=12, bold=True,
                  color="GOLD", bg_hex="BLACK",
                  align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    # Data rows with alternating background
    for ri, (label, col) in enumerate(summary_metrics, start=1):
        bkfc_val = safe_float(data["match_bkfc"][col])
        opp_val  = safe_float(data["match_opp"][col])
        bg = "ROW_ALT" if ri % 2 == 0 else "WHITE"

        _set_cell(tbl.cell(ri, 0), label,
                  size=11, bg_hex=bg)
        _set_cell(tbl.cell(ri, 1), f"{bkfc_val:.2g}",
                  size=11, bold=True, bg_hex=bg,
                  align=PP_ALIGN.CENTER)
        _set_cell(tbl.cell(ri, 2), f"{opp_val:.2g}",
                  size=11, bold=True, bg_hex=bg,
                  align=PP_ALIGN.CENTER)


# ── SLIDE 3: STRATEGIC INSIGHTS ──────────────────────────────────────────────
def add_insights_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, "STRATEGIC MATCH INSIGHTS",
                      "Key statistical deviations from season baseline")
    _add_slide_footer(slide, data)

    stats_list = [
        {
            "label":  s["label"],
            "match":  safe_float(data["match_bkfc"][s["col"]]),
            "season": safe_float(data["bkfc_season_avg"][s["col"]]),
        }
        for s in STATS
    ]
    insights = generate_insights(stats_list)

    y = 1.1
    for text in insights:
        # Card background
        card = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(12.13), Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = get_rgb(COLORS["DARK_GRAY"])
        card.line.color.rgb = get_rgb(COLORS["GOLD"])
        card.line.width = Pt(0.75)

        _add_text(slide, f"▪  {text}",
                  0.8, y + 0.07, 11.9, 0.62,
                  size=11, bold=False, color="WHITE")
        y += 0.9


# ── SLIDE 4: HEAD-TO-HEAD COMPARISON ────────────────────────────────────────
def add_h2h_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, "HEAD-TO-HEAD COMPARISON",
                      f"BKFC vs {data['opponent_name'].upper()} — all key metrics this match")
    _add_slide_footer(slide, data)

    bkfc_vals = np.array([safe_float(data["match_bkfc"][s["col"]]) for s in H2H_STATS])
    opp_vals  = np.array([safe_float(data["match_opp"][s["col"]]) for s in H2H_STATS])
    labels    = [s["label"] for s in H2H_STATS]

    chart = _h_bar_chart(bkfc_vals, opp_vals, labels,
                         f"BKFC vs {data['opponent_name']}")
    slide.shapes.add_picture(chart, Inches(0.5), Inches(1.1),
                             Inches(12.33), Inches(5.9))


# ── SLIDE 5: ATTACKING PHASE BREAKDOWN ──────────────────────────────────────
def add_attacking_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, "ATTACKING PHASE BREAKDOWN",
                      "Positional attacks · Counters · Set pieces · Ball progression")
    _add_slide_footer(slide, data)

    # Left chart: this match BKFC vs Opp
    match_vals = [safe_float(data["match_bkfc"][s["col"]]) for s in ATTACK_STATS]
    opp_vals   = [safe_float(data["match_opp"][s["col"]]) for s in ATTACK_STATS]
    labels     = [s["label"] for s in ATTACK_STATS]

    # We build a grouped bar for this match
    n = len(labels)
    x = np.arange(n)
    bar_w = 0.35

    fig, ax = plt.subplots(figsize=(6.0, 4.5), dpi=170)
    fig.patch.set_facecolor("#" + COLORS["WHITE"])
    ax.set_facecolor("#" + COLORS["WHITE"])

    b1 = ax.bar(x - bar_w/2, match_vals, bar_w,
                color="#" + COLORS["GOLD"], label="BKFC", edgecolor="none")
    b2 = ax.bar(x + bar_w/2, opp_vals, bar_w,
                color="#" + COLORS["DARK_GRAY"], label=data["opponent_name"],
                edgecolor="none")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=6.5, rotation=30, ha="right",
                       color="#" + COLORS["BLACK"])
    ax.tick_params(axis="y", colors="#" + COLORS["SILVER"], labelsize=7)
    ax.set_title("THIS MATCH — ATTACKING OUTPUT", fontsize=9,
                 fontweight="bold", pad=8, color="#" + COLORS["BLACK"])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.spines["bottom"].set_color("#" + COLORS["SILVER"])
    ax.yaxis.grid(True, linestyle="--", alpha=0.4,
                  color="#" + COLORS["GRID"], zorder=0)
    ax.set_axisbelow(True)
    ax.legend(fontsize=7.5, frameon=False)

    for bar in list(b1) + list(b2):
        h = bar.get_height()
        ax.annotate(f"{h:.2g}",
                    xy=(bar.get_x() + bar.get_width()/2, h),
                    xytext=(0, 3), textcoords="offset points",
                    ha="center", va="bottom", fontsize=7, fontweight="bold",
                    color="#" + COLORS["BLACK"])

    plt.tight_layout()
    buf1 = io.BytesIO()
    plt.savefig(buf1, format="png", facecolor=fig.get_facecolor(), edgecolor="none")
    buf1.seek(0)
    plt.close(fig)

    # Right chart: BKFC season avg vs all-opp avg
    season_vals = [safe_float(data["bkfc_season_avg"][s["col"]]) for s in ATTACK_STATS]
    league_vals = [safe_float(data["all_opp_avg"][s["col"]]) for s in ATTACK_STATS]

    buf2 = _bar_chart(
        season_vals + league_vals,
        [""] * n + [""] * n,   # labels suppressed; handled by positions
        [COLORS["GOLD"]] * n + [COLORS["SILVER"]] * n,
        "SEASON BASELINE — ATTACKING AVG"
    )

    # Simpler season comparison chart
    fig2, ax2 = plt.subplots(figsize=(5.6, 4.5), dpi=170)
    fig2.patch.set_facecolor("#" + COLORS["WHITE"])
    ax2.set_facecolor("#" + COLORS["WHITE"])
    b3 = ax2.bar(x - bar_w/2, season_vals, bar_w,
                 color="#" + COLORS["GOLD"], label="BKFC Avg", edgecolor="none")
    b4 = ax2.bar(x + bar_w/2, league_vals, bar_w,
                 color="#" + COLORS["SILVER"], label="Opp Avg", edgecolor="none")
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels, fontsize=6.5, rotation=30, ha="right",
                        color="#" + COLORS["BLACK"])
    ax2.tick_params(axis="y", colors="#" + COLORS["SILVER"], labelsize=7)
    ax2.set_title("SEASON BASELINE — ATTACKING AVG", fontsize=9,
                  fontweight="bold", pad=8, color="#" + COLORS["BLACK"])
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)
    ax2.spines["left"].set_visible(False)
    ax2.spines["bottom"].set_color("#" + COLORS["SILVER"])
    ax2.yaxis.grid(True, linestyle="--", alpha=0.4,
                   color="#" + COLORS["GRID"], zorder=0)
    ax2.set_axisbelow(True)
    ax2.legend(fontsize=7.5, frameon=False)
    plt.tight_layout()
    buf2 = io.BytesIO()
    plt.savefig(buf2, format="png", facecolor=fig2.get_facecolor(), edgecolor="none")
    buf2.seek(0)
    plt.close(fig2)

    slide.shapes.add_picture(buf1, Inches(0.3), Inches(1.1), Inches(6.5), Inches(5.5))
    slide.shapes.add_picture(buf2, Inches(6.9), Inches(1.1), Inches(6.1), Inches(5.5))

    # Vertical divider
    div = slide.shapes.add_shape(
        1, Inches(6.66), Inches(1.2), Inches(0.02), Inches(5.4)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = get_rgb(COLORS["GRID"])
    div.line.fill.background()


# ── SLIDE 6: DEFENSIVE PHASE BREAKDOWN ──────────────────────────────────────
def add_defensive_slide(prs, data):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, "DEFENSIVE PHASE BREAKDOWN",
                      "Duels · Press intensity (PPDA) · Interceptions · Clearances")
    _add_slide_footer(slide, data)

    n = len(DEFENSE_STATS)
    x = np.arange(n)
    bar_w = 0.35
    labels = [s["label"] for s in DEFENSE_STATS]

    def _def_chart(bkfc_v, comp_v, comp_label, title):
        fig, ax = plt.subplots(figsize=(5.6, 4.5), dpi=170)
        fig.patch.set_facecolor("#" + COLORS["WHITE"])
        ax.set_facecolor("#" + COLORS["WHITE"])
        b1 = ax.bar(x - bar_w/2, bkfc_v, bar_w,
                    color="#" + COLORS["GOLD"], label="BKFC", edgecolor="none")
        b2 = ax.bar(x + bar_w/2, comp_v, bar_w,
                    color="#" + COLORS["DARK_GRAY"], label=comp_label,
                    edgecolor="none")
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=6.5, rotation=30, ha="right",
                           color="#" + COLORS["BLACK"])
        ax.tick_params(axis="y", colors="#" + COLORS["SILVER"], labelsize=7)
        ax.set_title(title, fontsize=8.5, fontweight="bold",
                     pad=8, color="#" + COLORS["BLACK"])
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_visible(False)
        ax.spines["bottom"].set_color("#" + COLORS["SILVER"])
        ax.yaxis.grid(True, linestyle="--", alpha=0.4,
                      color="#" + COLORS["GRID"], zorder=0)
        ax.set_axisbelow(True)
        ax.legend(fontsize=7.5, frameon=False)
        for bar in list(b1) + list(b2):
            h = bar.get_height()
            ax.annotate(f"{h:.2g}",
                        xy=(bar.get_x() + bar.get_width()/2, h),
                        xytext=(0, 3), textcoords="offset points",
                        ha="center", va="bottom", fontsize=7,
                        fontweight="bold", color="#" + COLORS["BLACK"])
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format="png",
                    facecolor=fig.get_facecolor(), edgecolor="none")
        buf.seek(0)
        plt.close(fig)
        return buf

    match_bkfc_v = [safe_float(data["match_bkfc"][s["col"]]) for s in DEFENSE_STATS]
    match_opp_v  = [safe_float(data["match_opp"][s["col"]]) for s in DEFENSE_STATS]
    season_bkfc_v = [safe_float(data["bkfc_season_avg"][s["col"]]) for s in DEFENSE_STATS]
    season_avg_v  = [safe_float(data["all_opp_avg"][s["col"]]) for s in DEFENSE_STATS]

    buf1 = _def_chart(match_bkfc_v, match_opp_v,
                      data["opponent_name"],
                      "THIS MATCH — DEFENSIVE OUTPUT")
    buf2 = _def_chart(season_bkfc_v, season_avg_v,
                      "Opp Avg",
                      "SEASON BASELINE — DEFENSIVE AVG")

    slide.shapes.add_picture(buf1, Inches(0.3), Inches(1.1), Inches(6.5), Inches(5.5))
    slide.shapes.add_picture(buf2, Inches(6.9), Inches(1.1), Inches(6.1), Inches(5.5))

    div = slide.shapes.add_shape(
        1, Inches(6.66), Inches(1.2), Inches(0.02), Inches(5.4)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = get_rgb(COLORS["GRID"])
    div.line.fill.background()


# ── SLIDE 7+: PER-STAT DEEP DIVE ────────────────────────────────────────────
def add_stat_slide(prs, data, label, col):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    _add_slide_header(slide, label.upper())
    _add_slide_footer(slide, data)

    # Sub-labels
    _add_text(slide, "THIS MATCH",
              0.5, 0.92, 6.0, 0.28,
              size=9.5, bold=True, color="SILVER",
              align=PP_ALIGN.CENTER)
    _add_text(slide, "SEASON CONTEXT",
              7.0, 0.92, 5.8, 0.28,
              size=9.5, bold=True, color="SILVER",
              align=PP_ALIGN.CENTER)

    # Vertical divider
    div = slide.shapes.add_shape(
        1, Inches(6.66), Inches(1.1), Inches(0.02), Inches(5.5)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = get_rgb(COLORS["GRID"])
    div.line.fill.background()

    # Chart 1 — This match
    match_vals   = [safe_float(data["match_bkfc"][col]),
                    safe_float(data["match_opp"][col])]
    match_labels = ["BKFC", data["opponent_name"]]
    match_colors = [COLORS["GOLD"], COLORS["DARK_GRAY"]]
    buf1 = _bar_chart(match_vals, match_labels, match_colors,
                      "Match Comparison")

    # Chart 2 — Season context
    season_vals   = [
        safe_float(data["bkfc_season_avg"][col]),
        safe_float(data["all_opp_avg"][col]),
        safe_float(data["opp_season_avg"][col]),
    ]
    season_labels = ["BKFC Avg", "Opp Avg", f"{data['opponent_name']} Avg"]
    season_colors = [COLORS["GOLD"], COLORS["SILVER"], COLORS["DARK_GRAY"]]
    buf2 = _bar_chart(season_vals, season_labels, season_colors,
                      "Historical Context")

    slide.shapes.add_picture(buf1, Inches(0.5), Inches(1.3), Inches(5.7), Inches(5.4))
    slide.shapes.add_picture(buf2, Inches(7.0), Inches(1.3), Inches(5.7), Inches(5.4))


# ── MAIN EXPORT PIPELINE ─────────────────────────────────────────────────────
def generate_report(data) -> io.BytesIO:
    """
    Compile all slides into a PowerPoint BytesIO stream.
    Slide order:
        1. Title
        2. Match Performance Overview (summary table)
        3. Strategic Insights
        4. Head-to-Head Comparison
        5. Attacking Phase Breakdown
        6. Defensive Phase Breakdown
        7+. Per-stat deep dives
    """
    prs = create_prs()

    add_title_slide(prs, data)
    add_summary_slide(prs, data)
    add_insights_slide(prs, data)
    add_h2h_slide(prs, data)
    add_attacking_slide(prs, data)
    add_defensive_slide(prs, data)

    for stat in STATS:
        add_stat_slide(prs, data, stat["label"], stat["col"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
